"""Deck engine — generates a .pptx presentation from a completed UX audit.

Every slide is drawn from scratch (shapes, text boxes, icon chips) on top of the
adesso template's blank layout. We keep only three things from the corporate
template: the adesso logo (extracted as an image asset), the adesso brand
colors (blue / grey), and Fira Sans — the one adesso corporate font that's
freely licensed for use outside the marketing department (see the styleguide
PDF, pages 24-27). Everything else — layout, icons, severity/verification
color coding — is a purpose-built template for this tool, sharing its color
language with the Streamlit app's own UI (styles.py) so the deck feels like a
continuation of the on-screen review, not a generic export.

PDF conversion is skipped on Streamlit Cloud (no LibreOffice).
"""

import os
import pathlib
import re
import sys
import tempfile
import zipfile
from datetime import datetime

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

_SCRIPTS_DIR = pathlib.Path(__file__).parent / "skills" / "ux-audit-deck-skill" / "scripts"
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from pptx_helpers import add_slide, dedup_zip  # noqa: E402

ASSETS_DIR = pathlib.Path(__file__).parent / "skills" / "ux-audit-deck-skill" / "assets"
TEMPLATE_PATH = ASSETS_DIR / "adesso_template.pptx"
LOGO_PATH = ASSETS_DIR / "adesso_logo.png"
LOGO_WHITE_PATH = ASSETS_DIR / "adesso_logo_white.png"

# Blank layout in the adesso template — we draw every slide on this, keeping
# only the template's page geometry/theme, not its placeholders.
# See skills/ux-audit-deck-skill/references/adesso_layouts.md.
LAYOUT_BLANK = 32

# ---------------------------------------------------------------------------
# Brand system
# ---------------------------------------------------------------------------
# Core adesso brand colors (styleguide pp.24-25).
BLUE = RGBColor(0x00, 0x6E, 0xC7)
BLUE_DARK = RGBColor(0x00, 0x53, 0x96)
BLUE_TINT = RGBColor(0xEA, 0xF3, 0xFC)
GREY = RGBColor(0x88, 0x7D, 0x75)
BG = RGBColor(0xF3, 0xF2, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1C, 0x1C, 0x1A)
BODY_TEXT = RGBColor(0x3D, 0x3D, 0x3A)
BORDER = RGBColor(0xE2, 0xDF, 0xD9)

# Severity + verification colors — deliberately reused from the app's own
# styles.py so the deck reads as a continuation of the on-screen review.
SEVERITY = {
    "major": (RGBColor(0xA3, 0x27, 0x1D), RGBColor(0xFE, 0xF2, 0xF2), "!"),
    "mid": (RGBColor(0xA3, 0x67, 0x1D), RGBColor(0xFF, 0xFB, 0xEB), "!"),
    "minor": (RGBColor(0x8A, 0x7A, 0x12), RGBColor(0xFE, 0xFC, 0xE8), "i"),
    "open": (RGBColor(0x2F, 0x5F, 0xA3), RGBColor(0xEF, 0xF6, 0xFF), "?"),
}
VERIFICATION = {
    "Verified": (RGBColor(0x15, 0x7A, 0x45), RGBColor(0xEC, 0xFD, 0xF3), "\u2713 VERIFIED"),
    "Unverified": (RGBColor(0x71, 0x71, 0x7A), RGBColor(0xF4, 0xF4, 0xF5), "? UNVERIFIED"),
    "Outdated?": (RGBColor(0xB4, 0x53, 0x09), RGBColor(0xFF, 0xFB, 0xEB), "\u26A0 OUTDATED?"),
}

FONT_HEAD = "Fira Sans Condensed"
FONT_BODY = "Fira Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)

VERIFICATION_TAG_RE = re.compile(r"\[(Verified|Unverified|Outdated\?)\]")

# Language-specific labels
LABELS = {
    "English": {
        "report_title": "UX Audit Report",
        "overview": "Overview",
        "major_findings": "Major Findings",
        "mid_findings": "Mid Findings",
        "minor_findings": "Minor Findings",
        "open_findings": "Open Questions",
        "benchmarks": "Benchmarks",
        "market_benchmark": "Market Benchmark",
        "competitor_benchmark": "Competitor Benchmark",
        "competitor_gallery": "Products Reviewed",
        "closing": "Thank you",
        "closing_body": "Next steps: review findings with the team, prioritize by severity, and schedule validation research for the top items.",
        "pages_reviewed": "Pages reviewed",
        "target_users": "Target users",
        "mode": "Mode",
        "severity_tally": "Severity tally",
        "benchmark_scope": "Benchmark scope",
        "finding": "Finding",
        "suggestions": "Suggestions",
        "solution": "Solution",
        "heuristic": "Heuristic",
        "no_screenshot": "No screenshot provided",
        "screens_reviewed": "Screens reviewed",
        "legend": "Verified via web search \u00b7 Unverified — drafted, not confirmed \u00b7 Possibly outdated",
        "not_specified": "Not specified",
    },
    "Türkçe": {
        "report_title": "UX Denetim Raporu",
        "overview": "Genel Bakış",
        "major_findings": "Önemli Tespitler",
        "mid_findings": "Orta Tespitler",
        "minor_findings": "Hafif Tespitler",
        "open_findings": "Açık Sorular",
        "benchmarks": "Kıyaslamalar",
        "market_benchmark": "Pazar Kıyaslaması",
        "competitor_benchmark": "Rakip Kıyaslaması",
        "competitor_gallery": "İncelenen Ürünler",
        "closing": "Teşekkürler",
        "closing_body": "Sonraki adımlar: tespitleri ekiple gözden geçirin, önem derecesine göre önceliklendirin ve en önemli maddeleri doğrulama araştırmasına planlayın.",
        "pages_reviewed": "İncelenen sayfalar",
        "target_users": "Hedef kullanıcılar",
        "mode": "Mod",
        "severity_tally": "Tespit dağılımı",
        "benchmark_scope": "Kıyaslama kapsamı",
        "finding": "Tespit",
        "suggestions": "Öneriler",
        "solution": "Çözüm",
        "heuristic": "Sezgisel kural",
        "no_screenshot": "Ekran görüntüsü yok",
        "screens_reviewed": "İncelenen ekranlar",
        "legend": "Web aramasıyla doğrulandı \u00b7 Doğrulanmadı — eğitim verisinden yazıldı \u00b7 Güncelliğini yitirmiş olabilir",
        "not_specified": "Belirtilmedi",
    },
    "Deutsch": {
        "report_title": "UX-Audit-Bericht",
        "overview": "Überblick",
        "major_findings": "Wichtige Befunde",
        "mid_findings": "Mittlere Befunde",
        "minor_findings": "Geringfügige Befunde",
        "open_findings": "Offene Fragen",
        "benchmarks": "Benchmarks",
        "market_benchmark": "Marktvergleich",
        "competitor_benchmark": "Wettbewerbsvergleich",
        "competitor_gallery": "Geprüfte Produkte",
        "closing": "Vielen Dank",
        "closing_body": "Nächste Schritte: Befunde im Team besprechen, nach Schweregrad priorisieren und Validierungsforschung für die wichtigsten Punkte planen.",
        "pages_reviewed": "Geprüfte Seiten",
        "target_users": "Zielnutzer",
        "mode": "Modus",
        "severity_tally": "Befundverteilung",
        "benchmark_scope": "Benchmark-Umfang",
        "finding": "Befund",
        "suggestions": "Vorschläge",
        "solution": "Lösung",
        "heuristic": "Heuristik",
        "no_screenshot": "Kein Screenshot vorhanden",
        "screens_reviewed": "Geprüfte Bildschirme",
        "legend": "Per Websuche bestätigt \u00b7 Unbestätigt — aus Trainingswissen verfasst \u00b7 Möglicherweise veraltet",
        "not_specified": "Nicht angegeben",
    },
}

SEVERITY_SECTION_KEYS = [
    ("major", "major_findings"),
    ("mid", "mid_findings"),
    ("minor", "minor_findings"),
    ("open", "open_findings"),
]


def _get_labels(language: str) -> dict:
    return LABELS.get(language, LABELS["English"])


def _fix_template(template_path: str) -> str:
    """Fix the content type in the template if needed."""
    fixed_path = os.path.join(tempfile.gettempdir(), "adesso_template_fixed.pptx")
    with zipfile.ZipFile(template_path, "r") as zin:
        with zipfile.ZipFile(fixed_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = data.replace(
                        b"presentationml.template.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)
    return fixed_path


def _open_template():
    from pptx import Presentation
    template = str(TEMPLATE_PATH)
    try:
        return Presentation(template)
    except Exception:
        return Presentation(_fix_template(template))


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1), shape_type=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    sp.shadow.inherit = False
    return sp


def _textbox(slide, x, y, w, h, anchor=None, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb


def _set_run(run, text, size, color, bold=False, italic=False, font=FONT_BODY, spacing=None):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def _text(slide, x, y, w, h, text, size, color, bold=False, italic=False, font=FONT_BODY,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, line_spacing=None, wrap=True):
    tb = _textbox(slide, x, y, w, h, anchor=anchor, wrap=wrap)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    _set_run(p.add_run(), text, size, color, bold, italic, font, spacing)
    return tb


def _tagged_paragraph(tf, text, size, color, first=False, align=PP_ALIGN.LEFT,
                       line_spacing=1.18, space_after=8, bullet=True):
    """Add a paragraph, rendering inline [Verified]/[Unverified]/[Outdated?]
    tags as small colored bold badges within the text flow — this is how
    benchmark verification status actually reaches the deck.
    """
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)

    prefix = "\u2022  " if bullet else ""
    text = prefix + text

    pos = 0
    for m in VERIFICATION_TAG_RE.finditer(text):
        if m.start() > pos:
            _set_run(p.add_run(), text[pos:m.start()], size, color, font=FONT_BODY)
        tag_color, tag_bg, tag_label = VERIFICATION[m.group(1)]
        run = p.add_run()
        _set_run(run, " " + tag_label + " ", size - 1.5, tag_color, bold=True, font=FONT_BODY)
        pos = m.end()
    if pos < len(text):
        _set_run(p.add_run(), text[pos:], size, color, font=FONT_BODY)
    return p


def _icon_chip(slide, cx, cy, diameter, glyph, color, bg=None, filled=False):
    """A thin-outline circular icon chip with a centered glyph — echoes the
    adesso pictogram style (minimal contour, single color per composition)."""
    x, y = cx - diameter // 2, cy - diameter // 2
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    if filled:
        sp.fill.solid()
        sp.fill.fore_color.rgb = bg or color
        sp.line.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = bg or WHITE
        sp.line.color.rgb = color
        sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    glyph_color = WHITE if filled else color
    _set_run(p.add_run(), glyph, Emu(diameter).inches * 40, glyph_color, bold=True, font=FONT_BODY)
    return sp


def _badge(slide, x, y, text, fg, bg, size=10):
    """A small pill badge (severity tag)."""
    w = Inches(0.09 * len(text) + 0.4)
    h = Inches(0.26)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = 0.5
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = bg
    sp.line.color.rgb = fg
    sp.line.width = Pt(1)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), text, size, fg, bold=True, font=FONT_BODY, spacing=0.4)
    return sp, w


def _accent_bar(slide, color, y=0, height=Inches(0.09)):
    return _rect(slide, 0, y, SLIDE_W, height, fill=color)


def _logo(slide, on_dark=False):
    path = str(LOGO_WHITE_PATH if on_dark else LOGO_PATH)
    w = Inches(0.85)
    h = Inches(0.34)
    slide.shapes.add_picture(path, SLIDE_W - MARGIN - w, SLIDE_H - Inches(0.5), w, h)


def _footer(slide, product_name, page_num, total_pages, on_dark=False):
    color = RGBColor(0xC7, 0xDF, 0xF5) if on_dark else GREY
    _text(slide, MARGIN, SLIDE_H - Inches(0.46), Inches(6), Inches(0.3),
          product_name, 8.5, color, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, SLIDE_W - Inches(2.65), SLIDE_H - Inches(0.46), Inches(0.9), Inches(0.3),
          f"{page_num:02d} / {total_pages:02d}", 8.5, color, font=FONT_BODY,
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _fitted_picture(slide, image_bytes, x, y, w, h, ext="png"):
    """Insert an image scaled to fit within (w, h), centered, preserving
    aspect ratio — never crops, never overflows."""
    from PIL import Image as PILImage

    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name
    try:
        with PILImage.open(tmp_path) as img:
            img_w, img_h = img.size
        scale = min(w / img_w, h / img_h)
        new_w, new_h = int(img_w * scale), int(img_h * scale)
        new_x = x + (w - new_w) // 2
        new_y = y + (h - new_h) // 2
        return slide.shapes.add_picture(tmp_path, new_x, new_y, new_w, new_h)
    finally:
        os.unlink(tmp_path)


def _framed_picture(slide, image_bytes, x, y, w, h, ext="png", border=BORDER):
    """A screenshot in a bordered white card — the visual-aid treatment used
    throughout the deck (findings, cover, overview, benchmarks)."""
    _rect(slide, x, y, w, h, fill=WHITE, line=border, line_w=Pt(1))
    pad = Inches(0.08)
    _fitted_picture(slide, image_bytes, x + pad, y + pad, w - 2 * pad, h - 2 * pad, ext=ext)


def _truncate_for_slide(text: str, max_words: int) -> str:
    """Truncate text to approximately max_words words, cleaning markdown."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"^-\s+", "\u2022 ", text, flags=re.MULTILINE)
    text = re.sub(r"^\d+\.\s+", "", text, flags=re.MULTILINE)
    words = text.split()
    if len(words) <= max_words:
        return text
    return " ".join(words[:max_words]) + "\u2026"


def _ext_from_mime(mime_type) -> str:
    if not mime_type:
        return "png"
    return mime_type.split("/")[-1].replace("jpeg", "jpg")


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_cover(prs, labels, product_name, language, mode, hero_shot):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)
    _rect(slide, 0, SLIDE_H - Inches(1.4), Inches(4.2), Inches(1.4), fill=BLUE_DARK)

    pill = _rect(slide, MARGIN, Inches(0.55), Inches(1.7), Inches(0.36), fill=None,
                  line=WHITE, line_w=Pt(1.25), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        pill.adjustments[0] = 0.5
    except Exception:
        pass
    tf = pill.text_frame
    tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), language.upper(), 10, WHITE, bold=True, font=FONT_BODY, spacing=0.6)

    title_w = Inches(7.6) if hero_shot else Inches(11.5)
    _text(slide, MARGIN, Inches(2.9), title_w, Inches(0.5), labels["report_title"].upper(),
          13, RGBColor(0xBF, 0xDD, 0xF7), bold=True, font=FONT_HEAD, spacing=1.2)
    _text(slide, MARGIN - Inches(0.05), Inches(3.35), title_w, Inches(2.2), product_name,
          40, WHITE, bold=True, font=FONT_HEAD, line_spacing=1.02)
    mode_labels = {"screenshots": "Screenshots", "urls": "URLs", "hybrid": "Hybrid"}
    _text(slide, MARGIN, Inches(5.7), title_w, Inches(0.4),
          f"{datetime.now().strftime('%B %Y')}  \u00b7  {mode_labels.get(mode, mode)} mode",
          13, RGBColor(0xD8, 0xEA, 0xFA), font=FONT_BODY)

    if hero_shot:
        fw, fh = Inches(4.6), Inches(4.9)
        fx, fy = SLIDE_W - MARGIN - fw, Inches(1.55)
        _rect(slide, fx + Inches(0.12), fy + Inches(0.12), fw, fh, fill=BLUE_DARK)
        _framed_picture(slide, hero_shot["data"], fx, fy, fw, fh,
                         ext=_ext_from_mime(hero_shot.get("mime_type")))

    _logo(slide, on_dark=True)
    return slide


def _build_overview(prs, labels, product_name, meta_items, tally, screenshots, page_num, total_pages):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    _accent_bar(slide, BLUE)
    _text(slide, MARGIN, Inches(0.4), Inches(6), Inches(0.32), "OVERVIEW", 12, BLUE,
          bold=True, font=FONT_HEAD, spacing=1.2)
    _text(slide, MARGIN, Inches(0.72), Inches(10), Inches(0.65), labels["overview"], 26,
          DARK, bold=True, font=FONT_HEAD)

    card_w, card_h = Inches(5.85), Inches(0.92)
    gap = Inches(0.2)
    x0, y0 = MARGIN, Inches(1.65)
    for i, (glyph, label, value) in enumerate(meta_items):
        col, row = i % 2, i // 2
        cx = x0 + col * (card_w + gap)
        cy = y0 + row * (card_h + gap)
        _rect(slide, cx, cy, card_w, card_h, fill=WHITE, line=BORDER, line_w=Pt(1),
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _icon_chip(slide, cx + Inches(0.55), cy + card_h // 2, Inches(0.52), glyph, BLUE)
        _text(slide, cx + Inches(0.95), cy + Inches(0.13), card_w - Inches(1.15), Inches(0.28),
              label.upper(), 9.5, GREY, bold=True, font=FONT_BODY, spacing=0.5)
        _text(slide, cx + Inches(0.95), cy + Inches(0.4), card_w - Inches(1.15), Inches(0.46),
              value, 13, DARK, font=FONT_BODY, line_spacing=1.05)

    n_rows = (len(meta_items) + 1) // 2
    y_after = y0 + n_rows * (card_h + gap) + Inches(0.06)

    if tally:
        _text(slide, MARGIN, y_after, Inches(4), Inches(0.28), labels["severity_tally"].upper(),
              9.5, GREY, bold=True, font=FONT_BODY, spacing=0.5)
        bx, by = MARGIN, y_after + Inches(0.3)
        for sev_key, label_key in SEVERITY_SECTION_KEYS:
            n = tally.get(sev_key, 0)
            if not n:
                continue
            fg, bg, _ = SEVERITY[sev_key]
            _, badge_w = _badge(slide, bx, by, f"{n} {labels[label_key]}", fg, bg, size=10)
            bx += badge_w + Inches(0.14)
        y_after = by + Inches(0.42)

    if screenshots:
        shown = screenshots[:5]
        _text(slide, MARGIN, y_after, Inches(6), Inches(0.28), labels["screens_reviewed"].upper(),
              9.5, GREY, bold=True, font=FONT_BODY, spacing=0.5)
        thumb_w, thumb_h = Inches(1.5), Inches(1.0)
        tx, ty = MARGIN, y_after + Inches(0.3)
        for s in shown:
            _framed_picture(slide, s["data"], tx, ty, thumb_w, thumb_h,
                             ext=_ext_from_mime(s.get("mime_type")))
            tx += thumb_w + Inches(0.14)
        remaining = len(screenshots) - len(shown)
        if remaining > 0:
            _text(slide, tx, ty, Inches(1), thumb_h, f"+{remaining}", 14, GREY,
                  font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)

    _logo(slide)
    _footer(slide, product_name, page_num, total_pages)
    return slide


def _build_section_divider(prs, product_name, subtitle, count_text, color, glyph, page_num, total_pages):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    block_w = Inches(4.6)
    _rect(slide, 0, 0, block_w, SLIDE_H, fill=color)
    _icon_chip(slide, block_w // 2, Inches(3.3), Inches(1.15), glyph, WHITE, bg=color, filled=False)

    tx = block_w + Inches(0.75)
    _text(slide, tx, Inches(2.95), Inches(7.2), Inches(0.35), subtitle.upper(), 12, color,
          bold=True, font=FONT_HEAD, spacing=1.2)
    _text(slide, tx - Inches(0.05), Inches(3.28), Inches(7.3), Inches(1.1), subtitle, 34,
          DARK, bold=True, font=FONT_HEAD)
    _text(slide, tx, Inches(4.2), Inches(6.5), Inches(0.4), count_text, 14, BODY_TEXT, font=FONT_BODY)

    _logo(slide)
    _footer(slide, product_name, page_num, total_pages)
    return slide


def _split_finding_body(body_text: str, labels: dict):
    """Pull Finding / Suggestion-Solution / Heuristic blocks out of a
    finding's Markdown body for individually-styled layout blocks."""
    blocks = []
    finding_match = re.search(
        r"\*\*(?:Finding|Tespit|Befund)\s*:?\*\*\s*(.+?)(?=\n\*\*|\Z)", body_text, re.DOTALL
    )
    main_text = finding_match.group(1).strip() if finding_match else body_text.strip().split("\n\n")[0]
    blocks.append((labels.get("finding", "Finding"), _truncate_for_slide(main_text, 55), "!"))

    sugg_match = re.search(
        r"\*\*(?:Suggestions|Öneriler|Vorschläge|Solution|Çözüm|Lösung)\s*:?\*\*\s*(.+?)(?=\n\*\*|\Z)",
        body_text, re.DOTALL,
    )
    if sugg_match:
        label_key = "solution" if re.search(r"Solution|Çözüm|Lösung", sugg_match.group(0)) else "suggestions"
        blocks.append((labels.get(label_key, "Solution"), _truncate_for_slide(sugg_match.group(1).strip(), 40), "\u2713"))
    return blocks


def _build_finding_slide(prs, labels, product_name, finding, shot, index, total, page_num, total_pages):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    fg, bg, glyph = SEVERITY.get(finding["severity"], SEVERITY["open"])
    _accent_bar(slide, fg)

    has_shot = shot is not None
    left_w = Inches(6.3) if has_shot else Inches(11.9)

    _badge(slide, MARGIN, Inches(0.5), finding["severity_raw"].upper(), fg, bg, size=10.5)
    _text(slide, SLIDE_W - MARGIN - Inches(1.4), Inches(0.5), Inches(1.4), Inches(0.3),
          f"{index:02d} / {total:02d}", 10, GREY, font=FONT_BODY, align=PP_ALIGN.RIGHT)

    title = finding["title"]
    _text(slide, MARGIN, Inches(1.0), left_w, Inches(1.0), title, 22, DARK, bold=True,
          font=FONT_HEAD, line_spacing=1.05)

    body_top = Inches(1.95)
    blocks = _split_finding_body(finding.get("body", ""), labels)
    tb = _textbox(slide, MARGIN, body_top, left_w, Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for label, text, glyph_b in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0 if p is tf.paragraphs[0] else 14)
        p.line_spacing = 1.15
        r1 = p.add_run()
        _set_run(r1, glyph_b + "  " + label.upper() + "\n", 10.5, fg, bold=True, font=FONT_BODY, spacing=0.4)
        r2 = p.add_run()
        _set_run(r2, text, 13, BODY_TEXT, font=FONT_BODY)

    if finding.get("heuristic"):
        heur_y = Inches(6.35)
        _rect(slide, MARGIN, heur_y, left_w, Inches(0.55), fill=BG, line=None,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(slide, MARGIN + Inches(0.18), heur_y, left_w - Inches(0.36), Inches(0.55),
              finding["heuristic"], 10.5, GREY, italic=True, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.05)

    if has_shot:
        fw, fh = Inches(5.6), Inches(5.35)
        fx, fy = SLIDE_W - MARGIN - fw, Inches(1.0)
        _framed_picture(slide, shot["data"], fx, fy, fw, fh, ext=_ext_from_mime(shot.get("mime_type")))
    else:
        fw, fh = Inches(0), Inches(0)

    _logo(slide)
    _footer(slide, product_name, page_num, total_pages)
    return slide


def _pack_benchmark_slides(benchmark_text: str, max_words_per_slide: int = 130, max_slides: int = 4):
    """Split the raw benchmark Markdown into paragraph-sized chunks that each
    fit comfortably on one slide, preserving verification tags for later
    inline badge rendering. Returns a list of plain-text chunks (tags intact).
    """
    text = re.sub(r"^##\s+.+\n?", "", benchmark_text.strip(), count=1)  # drop the "## Benchmarking..." heading
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    slides, current, current_words = [], [], 0
    for para in paragraphs:
        para_clean = re.sub(r"^#{2,3}\s+", "", para).strip()
        n_words = len(para_clean.split())
        if current and current_words + n_words > max_words_per_slide and len(slides) < max_slides - 1:
            slides.append("\n\n".join(current))
            current, current_words = [], 0
        current.append(para_clean)
        current_words += n_words
    if current:
        slides.append("\n\n".join(current))
    return slides[:max_slides] if slides else []


def _build_benchmark_slide(prs, labels, product_name, title, text_chunk, shot, show_legend, page_num, total_pages):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
    _accent_bar(slide, BLUE)

    has_shot = shot is not None
    left_w = Inches(6.6) if has_shot else Inches(11.9)

    _text(slide, MARGIN, Inches(0.42), left_w, Inches(0.35), title.upper(), 12, BLUE,
          bold=True, font=FONT_HEAD, spacing=1.0)

    body_top = Inches(0.9)
    body_h = Inches(5.5) if show_legend else Inches(5.85)
    tb = _textbox(slide, MARGIN, body_top, left_w, body_h)
    tf = tb.text_frame
    tf.word_wrap = True

    paragraphs = [p for p in text_chunk.split("\n\n") if p.strip()]
    first = True
    for para in paragraphs:
        lines = [ln.strip("-• ").strip() for ln in para.split("\n") if ln.strip()]
        for line in lines:
            if not line:
                continue
            _tagged_paragraph(tf, line, 12.5, BODY_TEXT, first=first, line_spacing=1.22, space_after=10)
            first = False

    if show_legend:
        ly = SLIDE_H - Inches(0.95)
        _rect(slide, MARGIN, ly, left_w, Inches(0.02), fill=BORDER)
        _text(slide, MARGIN, ly + Inches(0.1), left_w, Inches(0.3), labels["legend"], 9, GREY,
              italic=True, font=FONT_BODY)

    if has_shot:
        fw, fh = Inches(5.6), Inches(6.0)
        fx, fy = SLIDE_W - MARGIN - fw, Inches(0.85)
        _framed_picture(slide, shot["data"], fx, fy, fw, fh, ext=_ext_from_mime(shot.get("mime_type")))

    _logo(slide)
    _footer(slide, product_name, page_num, total_pages)
    return slide


def _build_gallery_slide(prs, labels, product_name, title, items, page_num, total_pages):
    """A grid of screenshots with captions — used for the competitor gallery
    so benchmark competitors get a visual page of their own, not just prose."""
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)
    _accent_bar(slide, BLUE)
    _text(slide, MARGIN, Inches(0.4), Inches(8), Inches(0.32), title.upper(), 12, BLUE,
          bold=True, font=FONT_HEAD, spacing=1.2)

    n = len(items)
    cols = min(3, n) if n <= 6 else 4
    rows = (n + cols - 1) // cols
    gap = Inches(0.28)
    avail_w = SLIDE_W - 2 * MARGIN - gap * (cols - 1)
    avail_h = Inches(5.6) - gap * (rows - 1)
    card_w = avail_w / cols
    card_h = min(avail_h / rows, Inches(2.6))

    for i, (caption, shot) in enumerate(items):
        col, row = i % cols, i // cols
        x = MARGIN + col * (card_w + gap)
        y = Inches(1.1) + row * (card_h + gap)
        img_h = card_h - Inches(0.4)
        _framed_picture(slide, shot["data"], x, y, card_w, img_h, ext=_ext_from_mime(shot.get("mime_type")))
        _text(slide, x, y + img_h + Inches(0.04), card_w, Inches(0.3), caption, 10.5, DARK,
              bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER)

    _logo(slide)
    _footer(slide, product_name, page_num, total_pages)
    return slide


def _build_closing(prs, labels, product_name, page_num, total_pages):
    slide = add_slide(prs, LAYOUT_BLANK)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)
    _rect(slide, SLIDE_W - Inches(3.6), 0, Inches(3.6), Inches(1.3), fill=BLUE_DARK)
    _text(slide, MARGIN, Inches(2.9), Inches(11), Inches(1.0), labels["closing"], 40, WHITE,
          bold=True, font=FONT_HEAD)
    _text(slide, MARGIN, Inches(3.95), Inches(9.5), Inches(1.4), labels["closing_body"], 15,
          RGBColor(0xD8, 0xEA, 0xFA), font=FONT_BODY, line_spacing=1.3)
    _logo(slide, on_dark=True)
    _footer(slide, product_name, page_num, total_pages, on_dark=True)
    return slide


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_deck(
    findings: list[dict],
    product_name: str,
    language: str,
    mode: str,
    benchmark: str,
    screenshots: list[dict] | None = None,
    tally: dict[str, int] | None = None,
    target_users: str = "",
    pages_reviewed: str = "",
    benchmark_text: str = "",
    competitors: list[dict] | None = None,
    comp_screenshots: dict | None = None,
    progress_callback=None,
) -> bytes:
    """Generate a .pptx deck from structured findings.

    Args:
        findings: list of parsed finding dicts from audit_engine.parse_findings
            (each may include a `screenshot_index`, 1-based, matching upload order)
        product_name: name for the cover slide
        language: 'English', 'Türkçe', or 'Deutsch'
        mode: 'screenshots', 'urls', or 'hybrid'
        benchmark: 'none', 'market', 'competitors', or 'both'
        screenshots: list of dicts with 'name', 'data' (bytes), 'mime_type', in upload
            order — index N-1 corresponds to a finding's screenshot_index == N
        tally: severity tally dict
        target_users: for the overview slide
        pages_reviewed: for the overview slide
        benchmark_text: raw benchmark section Markdown from the audit, with inline
            [Verified]/[Unverified]/[Outdated?] tags per the skill's research protocol
        competitors: list of {"name", "url"} dicts, for captioning the competitor gallery
        comp_screenshots: dict[competitor_index -> list of {"name","data","mime_type"}]
        progress_callback: callable(current, total, label) for progress updates

    Returns:
        bytes of the generated .pptx file
    """
    labels = _get_labels(language)
    screenshots = screenshots or []
    competitors = competitors or []
    comp_screenshots = comp_screenshots or {}

    prs = _open_template()
    # The template ships with pre-populated demo slides — drop them all, we
    # build the entire deck fresh from the blank layout.
    while len(prs.slides) > 0:
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(list(xml_slides)[0])

    # Flatten competitor screenshots for the gallery + per-competitor lookup.
    gallery_items = []
    for i, comp in enumerate(competitors):
        shots = comp_screenshots.get(i) or []
        name = comp.get("name") or f"Competitor {i + 1}"
        for s in shots:
            gallery_items.append((name, s))

    benchmark_chunks = _pack_benchmark_slides(benchmark_text) if benchmark_text else []

    total_sections = sum(1 for sev, _ in SEVERITY_SECTION_KEYS if any(f["severity"] == sev for f in findings))
    total_slides = 2 + total_sections + len(findings) + 1  # cover + overview + sections + findings + closing
    if benchmark != "none" and benchmark_chunks:
        total_slides += 1 + len(benchmark_chunks)  # divider + content chunks
        if gallery_items:
            total_slides += 1

    page_num = 0

    def _progress(label):
        nonlocal page_num
        page_num += 1
        if progress_callback:
            progress_callback(page_num, total_slides, label)

    # --- 1. Cover ---
    hero_shot = screenshots[0] if screenshots else None
    _build_cover(prs, labels, product_name, language, mode, hero_shot)
    _progress("Cover")

    # --- 2. Overview ---
    mode_labels = {"screenshots": "Screenshots", "urls": "URLs", "hybrid": "Hybrid"}
    bench_labels = {"market": "Market", "competitors": "Competitors", "both": "Market + Competitors"}
    meta_items = [
        ("\u25A4", labels["pages_reviewed"], pages_reviewed or labels["not_specified"]),
        ("\u25CB", labels["target_users"], target_users or labels["not_specified"]),
        ("\u25C8", labels["mode"], mode_labels.get(mode, mode)),
    ]
    if benchmark != "none":
        meta_items.append(("\u2194", labels["benchmark_scope"], bench_labels.get(benchmark, benchmark)))
    _build_overview(prs, labels, product_name, meta_items, tally, screenshots, page_num + 1, total_slides)
    _progress("Overview")

    # --- 3-4. Section dividers + findings ---
    finding_counter = 0
    for sev_key, section_label_key in SEVERITY_SECTION_KEYS:
        sev_findings = [f for f in findings if f["severity"] == sev_key]
        if not sev_findings:
            continue

        fg, bg, glyph = SEVERITY[sev_key]
        count_text = f"{len(sev_findings)} finding" + ("s" if len(sev_findings) != 1 else "")
        _build_section_divider(prs, product_name, labels[section_label_key], count_text, fg, glyph,
                                page_num + 1, total_slides)
        _progress(labels[section_label_key])

        for f in sev_findings:
            finding_counter += 1
            shot = None
            shot_idx = f.get("screenshot_index")
            if shot_idx and 1 <= shot_idx <= len(screenshots):
                shot = screenshots[shot_idx - 1]
            _build_finding_slide(prs, labels, product_name, f, shot, finding_counter, len(findings),
                                  page_num + 1, total_slides)
            _progress(f["title"][:30])

    # --- 5. Benchmark section ---
    if benchmark != "none" and benchmark_chunks:
        bench_title = labels.get(
            "market_benchmark" if benchmark == "market" else "competitor_benchmark", "Benchmarks"
        )
        _build_section_divider(prs, product_name, labels["benchmarks"],
                                bench_title, BLUE, "\u2194", page_num + 1, total_slides)
        _progress("Benchmark divider")

        if gallery_items:
            _build_gallery_slide(prs, labels, product_name, labels["competitor_gallery"],
                                  gallery_items, page_num + 1, total_slides)
            _progress("Competitor gallery")

        # Pair the first couple of chunks with a competitor/market screenshot
        # when available, so benchmark slides carry visual evidence too.
        shot_pool = [s for _, s in gallery_items] or (screenshots[:1] if screenshots else [])
        for i, chunk in enumerate(benchmark_chunks):
            shot = shot_pool[i] if i < len(shot_pool) else None
            title = bench_title if i == 0 else f"{bench_title} \u2014 {i + 1}/{len(benchmark_chunks)}"
            _build_benchmark_slide(prs, labels, product_name, title, chunk, shot,
                                    show_legend=(i == 0), page_num=page_num + 1, total_pages=total_slides)
            _progress(f"Benchmark {i + 1}")

    # --- 6. Closing ---
    _build_closing(prs, labels, product_name, page_num + 1, total_slides)
    _progress("Closing")

    # Save to bytes
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_raw:
        prs.save(tmp_raw.name)
        raw_path = tmp_raw.name

    deduped_path = raw_path.replace(".pptx", "_clean.pptx")
    dedup_zip(raw_path, deduped_path)

    with open(deduped_path, "rb") as f:
        result = f.read()

    os.unlink(raw_path)
    os.unlink(deduped_path)

    return result
